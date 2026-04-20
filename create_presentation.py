#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint Sunum Oluşturucu
Markdown dosyasından yatırımcı sunumu oluşturur
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Yatırımcı sunumu PowerPoint dosyası oluştur"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Renkler (Modern, profesyonel tema)
    PRIMARY_COLOR = RGBColor(59, 130, 246)  # Mavi
    SECONDARY_COLOR = RGBColor(16, 185, 129)  # Yeşil
    DARK_COLOR = RGBColor(31, 41, 55)  # Koyu gri
    LIGHT_COLOR = RGBColor(243, 244, 246)  # Açık gri
    
    # SLAYT 1: KAPAK
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Başlık
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🎨 GÖRSEL DÖNÜŞTÜRÜCÜ"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    title_p.alignment = PP_ALIGN.CENTER
    
    # Alt başlık
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Profesyonel Görsel İşleme Platformu"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = DARK_COLOR
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Hibrit çözüm
    hybrid_box = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(8), Inches(0.5))
    hybrid_frame = hybrid_box.text_frame
    hybrid_frame.text = "Desktop + Web Hibrit Çözüm"
    hybrid_p = hybrid_frame.paragraphs[0]
    hybrid_p.font.size = Pt(22)
    hybrid_p.font.color.rgb = SECONDARY_COLOR
    hybrid_p.alignment = PP_ALIGN.CENTER
    
    # Tarih
    date_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.4))
    date_frame = date_box.text_frame
    date_frame.text = "Nisan 2026 • Türkiye"
    date_p = date_frame.paragraphs[0]
    date_p.font.size = Pt(18)
    date_p.font.color.rgb = DARK_COLOR
    date_p.alignment = PP_ALIGN.CENTER
    
    # SLAYT 2: PROBLEM
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "❌ PAZARIN ACILARI - 6 TEMEL SORUN"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    
    # Sorunlar - 2x3 grid
    problems = [
        ("🔒 GİZLİLİK", "Dosyalar sunuculara\nyükleniyor"),
        ("💰 YÜKSEK MALİYET", "Photoshop: $240/yıl\nCanva Pro: $120/yıl"),
        ("📉 KALİTE KAYBI", "Online araçlar\nsıkıştırıyor"),
        ("⏱️ YAVAŞ İŞLEM", "Sunucu yüklemesi +\nİşlem + İndirme"),
        ("🎨 SINIRLI ÖZELLİK", "Temel dönüşümler\nProfesyonel yok"),
        ("🌐 İNTERNET ZORUNLU", "Offline\nçalışmıyor")
    ]
    
    x_positions = [0.5, 5.25]
    y_positions = [1.5, 3.2, 4.9]
    
    idx = 0
    for row in range(3):
        for col in range(2):
            if idx < len(problems):
                x = x_positions[col]
                y = y_positions[row]
                
                box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4.5), Inches(1.3))
                frame = box.text_frame
                frame.word_wrap = True
                
                # Başlık
                p = frame.paragraphs[0]
                p.text = problems[idx][0]
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = DARK_COLOR
                
                # Açıklama
                p2 = frame.add_paragraph()
                p2.text = problems[idx][1]
                p2.font.size = Pt(14)
                p2.font.color.rgb = RGBColor(107, 114, 128)
                p2.level = 0
                
                idx += 1
    
    # SLAYT 3: ÇÖZÜM
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "✅ BİZİM ÇÖZÜMÜMÜZ"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = SECONDARY_COLOR
    
    # Desktop kutusu
    desktop_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.8))
    frame = desktop_box.text_frame
    frame.word_wrap = True
    
    p = frame.paragraphs[0]
    p.text = "🖥️ MASAÜSTÜ UYGULAMASI"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    features = ["Offline çalışır", "Dosyalar cihazınızda kalır", "Anında işleme (1-2 saniye)", "Sınırsız dosya boyutu"]
    for feature in features:
        p2 = frame.add_paragraph()
        p2.text = f"✓ {feature}"
        p2.font.size = Pt(16)
        p2.font.color.rgb = DARK_COLOR
        p2.level = 0
    
    # Web kutusu
    web_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1.8))
    frame2 = web_box.text_frame
    frame2.word_wrap = True
    
    p = frame2.paragraphs[0]
    p.text = "🌐 WEB UYGULAMASI"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    features2 = ["Platformdan bağımsız", "Kurulum gerektirmez", "Mobil uyumlu", "Güvenli API (otomatik dosya silme)"]
    for feature in features2:
        p2 = frame2.add_paragraph()
        p2.text = f"✓ {feature}"
        p2.font.size = Pt(16)
        p2.font.color.rgb = DARK_COLOR
        p2.level = 0
    
    # Alt mesaj
    msg_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.6))
    msg_frame = msg_box.text_frame
    msg_frame.text = "🎯 TEK PLATFORM - ÇOKLU ERİŞİM"
    msg_p = msg_frame.paragraphs[0]
    msg_p.font.size = Pt(22)
    msg_p.font.bold = True
    msg_p.font.color.rgb = SECONDARY_COLOR
    msg_p.alignment = PP_ALIGN.CENTER
    
    # SLAYT 4: ÜRÜN ÖZELLİKLERİ
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "🎨 ÜRÜN ÖZELLİKLERİ"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    
    features_list = [
        ("1️⃣ FORMAT DÖNÜŞTÜRME", "JPG • PNG • WebP • AVIF • TIFF • BMP"),
        ("2️⃣ AKILLI SIKIŞTIRMA", "Hedef dosya boyutu • Otomatik kalite • %70 küçültme"),
        ("3️⃣ PROFESYONEL RENK UZAYI", "CMYK ↔ RGB • ICC profil • Matbaa hazırlığı"),
        ("4️⃣ FİLİGRAN", "9 konum • Opacity • Gölge • Tile desteği"),
        ("5️⃣ TOPLU İŞLEM", "Free: 10 dosya • Premium: 100+ dosya")
    ]
    
    y_start = 1.2
    for i, (title, desc) in enumerate(features_list):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(y_start + i * 1.1), Inches(9), Inches(0.9))
        frame = box.text_frame
        frame.word_wrap = True
        
        p = frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_COLOR
        
        p2 = frame.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(107, 114, 128)
    
    # SLAYT 5: PAZAR BÜYÜKLÜĞÜ
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "📊 PAZAR BÜYÜKLÜĞÜ"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    
    # TAM
    tam_box = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(6), Inches(1.2))
    tam_frame = tam_box.text_frame
    tam_frame.text = "$2.3 MİLYAR / YIL"
    tam_p = tam_frame.paragraphs[0]
    tam_p.font.size = Pt(48)
    tam_p.font.bold = True
    tam_p.font.color.rgb = PRIMARY_COLOR
    tam_p.alignment = PP_ALIGN.CENTER
    
    tam_desc_box = slide.shapes.add_textbox(Inches(2), Inches(2.7), Inches(6), Inches(0.5))
    tam_desc_frame = tam_desc_box.text_frame
    tam_desc_frame.text = "Total Addressable Market (TAM)"
    tam_desc_p = tam_desc_frame.paragraphs[0]
    tam_desc_p.font.size = Pt(18)
    tam_desc_p.font.color.rgb = DARK_COLOR
    tam_desc_p.alignment = PP_ALIGN.CENTER
    
    # Segmentler
    segments = [
        "🎯 BİREYSEL KULLANICILAR: 1.2 milyar potansiyel",
        "🏢 KÜÇÜK İŞLETMELER: 350M SMB worldwide",
        "🎨 PROFESYONELLER: 23M grafik tasarımcı"
    ]
    
    y = 3.5
    for segment in segments:
        box = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.5))
        frame = box.text_frame
        p = frame.paragraphs[0]
        p.text = segment
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_COLOR
        y += 0.7
    
    # Büyüme
    growth_box = slide.shapes.add_textbox(Inches(2), Inches(6), Inches(6), Inches(0.6))
    growth_frame = growth_box.text_frame
    growth_frame.text = "📈 YILLIK BÜYÜME: %8.5 CAGR (2026-2031)"
    growth_p = growth_frame.paragraphs[0]
    growth_p.font.size = Pt(20)
    growth_p.font.bold = True
    growth_p.font.color.rgb = SECONDARY_COLOR
    growth_p.alignment = PP_ALIGN.CENTER
    
    # SLAYT 6: RAKİP ANALİZİ
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "⚔️ RAKIP KARŞILAŞTIRMA"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    
    # Tablo başlıkları ve veriler
    comparison_text = """
BİZ          Canva Pro    Adobe        Online Tools
────────────────────────────────────────────────────
✅ Offline   ❌ Online    ❌ Online    ❌ Online
✅ Gizli     ❌ Sunucu    ❌ Sunucu    ❌ Sunucu
✅ 1-2sn     ⚠️ 5-10sn    ⚠️ 5-10sn    ⚠️ 10sn+
✅ CMYK      ❌ Yok       ❌ Yok       ❌ Yok
✅ 100+      ⚠️ 10        ⚠️ 15        ⚠️ 5-10
$36/yıl      $120/yıl     $100/yıl     Ücretsiz
    """
    
    table_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(4))
    table_frame = table_box.text_frame
    table_frame.text = comparison_text
    table_p = table_frame.paragraphs[0]
    table_p.font.size = Pt(13)
    table_p.font.name = 'Consolas'
    table_p.font.color.rgb = DARK_COLOR
    
    # Farklar
    diff_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    diff_frame = diff_box.text_frame
    diff_frame.word_wrap = True
    
    p = diff_frame.paragraphs[0]
    p.text = "🏆 TEMEL FARKLILAŞMA:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    diffs = ["Hibrit Model (Desktop + Web)", "CMYK + ICC Profil", "Gizlilik (Yerel işleme)", "Maliyet Avantajı"]
    for diff in diffs:
        p2 = diff_frame.add_paragraph()
        p2.text = f"• {diff}"
        p2.font.size = Pt(14)
        p2.font.color.rgb = DARK_COLOR
    
    # SLAYT 7: TRAKSİYON
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "📈 TRAKSİYON (6 AY - BOOTSTRAP)"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    
    metrics = [
        ("Toplam Kayıt", "1,850 kullanıcı"),
        ("MAU (Active)", "1,200 kullanıcı"),
        ("Premium", "68 kullanıcı"),
        ("Conversion Rate", "%5.6"),
        ("MRR", "$204"),
        ("Retention (30-gün)", "%42")
    ]
    
    y = 1.5
    for metric, value in metrics:
        # Metrik adı
        label_box = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(4), Inches(0.5))
        label_frame = label_box.text_frame
        label_p = label_frame.paragraphs[0]
        label_p.text = metric
        label_p.font.size = Pt(18)
        label_p.font.color.rgb = DARK_COLOR
        
        # Değer
        value_box = slide.shapes.add_textbox(Inches(5.5), Inches(y), Inches(3), Inches(0.5))
        value_frame = value_box.text_frame
        value_p = value_frame.paragraphs[0]
        value_p.text = value
        value_p.font.size = Pt(22)
        value_p.font.bold = True
        value_p.font.color.rgb = PRIMARY_COLOR
        value_p.alignment = PP_ALIGN.RIGHT
        
        y += 0.8
    
    # Alt mesaj
    msg_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.8))
    msg_frame = msg_box.text_frame
    msg_frame.text = "🎯 ORGANIK BÜYÜME - ZERO MARKETING BUDGET"
    msg_p = msg_frame.paragraphs[0]
    msg_p.font.size = Pt(20)
    msg_p.font.bold = True
    msg_p.font.color.rgb = SECONDARY_COLOR
    msg_p.alignment = PP_ALIGN.CENTER
    
    # SLAYT 8: FİNANSAL PROJEKSİYON
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "💹 FİNANSAL PROJEKSİYON (3 YIL)"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    
    projections = [
        ("2026 (Yıl Sonu)", "25,000 MAU", "$60,000 ARR"),
        ("2027", "120,000 MAU", "$379,200 ARR"),
        ("2028", "450,000 MAU", "$1,614,000 ARR")
    ]
    
    y = 1.5
    for year, mau, arr in projections:
        # Yıl
        year_box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(2.5), Inches(0.6))
        year_frame = year_box.text_frame
        year_p = year_frame.paragraphs[0]
        year_p.text = year
        year_p.font.size = Pt(18)
        year_p.font.bold = True
        year_p.font.color.rgb = DARK_COLOR
        
        # MAU
        mau_box = slide.shapes.add_textbox(Inches(3.5), Inches(y), Inches(2.5), Inches(0.6))
        mau_frame = mau_box.text_frame
        mau_p = mau_frame.paragraphs[0]
        mau_p.text = mau
        mau_p.font.size = Pt(16)
        mau_p.font.color.rgb = PRIMARY_COLOR
        
        # ARR
        arr_box = slide.shapes.add_textbox(Inches(6.2), Inches(y), Inches(2.8), Inches(0.6))
        arr_frame = arr_box.text_frame
        arr_p = arr_frame.paragraphs[0]
        arr_p.text = arr
        arr_p.font.size = Pt(18)
        arr_p.font.bold = True
        arr_p.font.color.rgb = SECONDARY_COLOR
        arr_p.alignment = PP_ALIGN.RIGHT
        
        y += 1.2
    
    # Unit Economics
    unit_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(2))
    unit_frame = unit_box.text_frame
    unit_frame.word_wrap = True
    
    p = unit_frame.paragraphs[0]
    p.text = "🎯 UNIT ECONOMICS"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    unit_metrics = [
        "LTV:CAC Ratio: 8:1 🏆",
        "CAC: $30 | LTV: $240",
        "Payback Period: 4 ay",
        "EBITDA Margin: %54"
    ]
    
    for metric in unit_metrics:
        p2 = unit_frame.add_paragraph()
        p2.text = f"• {metric}"
        p2.font.size = Pt(16)
        p2.font.color.rgb = DARK_COLOR
    
    # SLAYT 9: YATIRIM TALEBİ
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "💰 YATIRIM TALEBİ"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = PRIMARY_COLOR
    
    # Ana tutar
    amount_box = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(6), Inches(1))
    amount_frame = amount_box.text_frame
    amount_frame.text = "$150,000"
    amount_p = amount_frame.paragraphs[0]
    amount_p.font.size = Pt(54)
    amount_p.font.bold = True
    amount_p.font.color.rgb = PRIMARY_COLOR
    amount_p.alignment = PP_ALIGN.CENTER
    
    # Seed Round
    seed_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(0.5))
    seed_frame = seed_box.text_frame
    seed_frame.text = "Seed Round • Equity: %15-20"
    seed_p = seed_frame.paragraphs[0]
    seed_p.font.size = Pt(20)
    seed_p.font.color.rgb = DARK_COLOR
    seed_p.alignment = PP_ALIGN.CENTER
    
    # Kullanım
    usage_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(2.5))
    usage_frame = usage_box.text_frame
    usage_frame.word_wrap = True
    
    p = usage_frame.paragraphs[0]
    p.text = "📊 FONLARIN KULLANIMI (18 Ay)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    usage_items = [
        "Personel (Ekip): $90,000 (60%)",
        "Pazarlama & Growth: $40,000 (27%)",
        "Infrastructure: $10,000 (7%)",
        "Legal & Admin: $10,000 (7%)"
    ]
    
    for item in usage_items:
        p2 = usage_frame.add_paragraph()
        p2.text = f"• {item}"
        p2.font.size = Pt(16)
        p2.font.color.rgb = DARK_COLOR
    
    # 18 ay hedef
    target_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(8.4), Inches(0.8))
    target_frame = target_box.text_frame
    target_frame.text = "🎯 18 Ay Hedef: 150,000 MAU • $216K ARR • Series A Hazır"
    target_p = target_frame.paragraphs[0]
    target_p.font.size = Pt(18)
    target_p.font.bold = True
    target_p.font.color.rgb = SECONDARY_COLOR
    target_p.alignment = PP_ALIGN.CENTER
    
    # SLAYT 10: KAPANIŞ
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Başlık
    closing_title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    closing_title_frame = closing_title_box.text_frame
    closing_title_frame.text = "PROFESYONEL GÖRSEL İŞLEMENİN\nDEMOKRATİKLEŞTİRİLMESİ"
    closing_title_p = closing_title_frame.paragraphs[0]
    closing_title_p.font.size = Pt(32)
    closing_title_p.font.bold = True
    closing_title_p.font.color.rgb = PRIMARY_COLOR
    closing_title_p.alignment = PP_ALIGN.CENTER
    
    # Özellikler
    features_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
    features_frame = features_box.text_frame
    features_frame.word_wrap = True
    
    closing_features = [
        "✓ Photoshop kalitesi",
        "✓ TinyPNG hızı",
        "✓ Offline gizliliği",
        "✓ 1/10 fiyatına"
    ]
    
    for feature in closing_features:
        p = features_frame.add_paragraph() if features_frame.paragraphs[0].text else features_frame.paragraphs[0]
        p.text = feature
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_COLOR
        p.alignment = PP_ALIGN.CENTER
    
    # Alt mesaj
    final_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1.2))
    final_frame = final_box.text_frame
    final_frame.word_wrap = True
    
    p = final_frame.paragraphs[0]
    p.text = "🤝 BİRLİKTE BÜYÜYELIM"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    p2 = final_frame.add_paragraph()
    p2.text = "Sorularınız için hazırız"
    p2.font.size = Pt(20)
    p2.font.color.rgb = DARK_COLOR
    p2.alignment = PP_ALIGN.CENTER
    
    # Dosyayı kaydet
    output_path = 'YATIRIMCI_SUNUMU.pptx'
    prs.save(output_path)
    print(f"✅ PowerPoint sunumu oluşturuldu: {output_path}")
    print(f"📊 Toplam {len(prs.slides)} slayt")
    
    return output_path

if __name__ == "__main__":
    create_presentation()
